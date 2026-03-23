import os
import sys
import re
from copy import deepcopy

from openpyxl import load_workbook
from pptx import Presentation

# =========================
# 기본 설정
# =========================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

TEMPLATE_PPT = os.path.join(BASE_DIR, "명찰 디자인.pptx")
OUTPUT_PPT = os.path.join(BASE_DIR, "완성_명찰.pptx")

NAME_COL = 1  # A열
PLACEHOLDERS = ["이름1", "이름2", "이름3", "이름4"]


# =========================
# 엑셀 파일 경로 받기
# =========================
def get_excel_path():
    if len(sys.argv) < 2:
        print("사용법: python creatMC.py 명단.xlsx")
        sys.exit(1)

    excel_path = sys.argv[1]

    if not os.path.exists(excel_path):
        print(f"[오류] 파일 없음: {excel_path}")
        sys.exit(1)

    return excel_path


# =========================
# 엑셀에서 이름 읽기
# =========================
def read_names_from_excel(path):
    wb = load_workbook(path, data_only=True)
    ws = wb.active

    names = []
    for row in ws.iter_rows(min_row=2, min_col=NAME_COL, max_col=NAME_COL, values_only=True):
        value = row[0]
        if value and str(value).strip():
            names.append(str(value).strip())

    return names


# =========================
# 슬라이드 복제
# =========================
def duplicate_slide(prs, source_slide):
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    rel_items = sorted(
        source_slide.part.rels.items(),
        key=lambda kv: int(re.sub(r"\D", "", kv[0]) or 0)
    )

    for rId, rel in rel_items:
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue

        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)

    return new_slide


# =========================
# 스타일 복사
# =========================
def copy_run_style(src_run, dst_run):
    if not src_run:
        return

    if src_run.font.name:
        dst_run.font.name = src_run.font.name
    if src_run.font.size:
        dst_run.font.size = src_run.font.size
    if src_run.font.bold is not None:
        dst_run.font.bold = src_run.font.bold
    if src_run.font.italic is not None:
        dst_run.font.italic = src_run.font.italic


# =========================
# 텍스트 치환 (핵심)
# =========================
def replace_paragraph(paragraph, mapping):
    full_text = "".join(run.text for run in paragraph.runs)
    if not full_text:
        return

    new_text = full_text
    for old, new in mapping.items():
        new_text = new_text.replace(old, new)

    if new_text == full_text:
        return

    first_run = paragraph.runs[0] if paragraph.runs else None

    # 기존 run 삭제
    for i in range(len(paragraph.runs) - 1, -1, -1):
        paragraph._p.remove(paragraph.runs[i]._r)

    # 새 run
    run = paragraph.add_run()
    run.text = new_text
    copy_run_style(first_run, run)


def replace_shape(shape, mapping):
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            replace_shape(s, mapping)

    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        replace_paragraph(paragraph, mapping)


def fill_slide(slide, names):
    mapping = {ph: "" for ph in PLACEHOLDERS}

    for i, name in enumerate(names):
        if i < 4:
            mapping[PLACEHOLDERS[i]] = name

    for shape in slide.shapes:
        replace_shape(shape, mapping)


# =========================
# 메인
# =========================
def main():
    excel_path = get_excel_path()

    if not os.path.exists(TEMPLATE_PPT):
        print(f"[오류] 템플릿 없음: {TEMPLATE_PPT}")
        return

    names = read_names_from_excel(excel_path)
    if not names:
        print("[오류] 이름 없음")
        return

    prs = Presentation(TEMPLATE_PPT)
    template_slide = prs.slides[0]

    chunks = [names[i:i+4] for i in range(0, len(names), 4)]

    # 슬라이드 복제
    for _ in range(len(chunks) - 1):
        duplicate_slide(prs, template_slide)

    # 이름 채우기
    for i, chunk in enumerate(chunks):
        fill_slide(prs.slides[i], chunk)

    prs.save(OUTPUT_PPT)
    print(f"완료: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
